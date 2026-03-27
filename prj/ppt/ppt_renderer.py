"""
ppt_renderer.py v4 (Academic Edition)
PPT rendering engine for paper2ppt
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree
from typing import List, Optional, Tuple
from ppt_style_analyzer import ThemeStyle
import re

W = Inches(13.33)
H = Inches(7.5)

ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
ACCENT_TEAL = RGBColor(0x1A, 0xBC, 0x9C)
ACADEMIC_BLUE = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x25, 0x3A)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)


def _rgb(h: str) -> RGBColor:
    if isinstance(h, RGBColor):
        return h
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def _is_rgb(val) -> bool:
    return isinstance(val, RGBColor)


def rect(slide, l, t, w, h, fill, line_color=None, line_width: float = 0.5):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    if _is_rgb(fill):
        s.fill.fore_color.rgb = fill
    else:
        s.fill.fore_color.rgb = _rgb(fill)
    if line_color:
        s.line.color.rgb = _rgb(line_color) if not _is_rgb(line_color) else line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def textbox(slide, text, l, t, w, h,
            font="Times New Roman", size=14, color="#000000",
            bold=False, italic=False,
            align=PP_ALIGN.LEFT, wrap=True,
            pad_l=0, pad_t=0, line_gap: int = 0):
    if not text:
        return None
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap

    bPr = tf._txBody.find(qn('a:bodyPr'))
    bPr.set('lIns', str(int(pad_l)))
    bPr.set('rIns', '0')
    bPr.set('tIns', str(int(pad_t)))
    bPr.set('bIns', '0')

    tf.paragraphs[0].clear()
    p = tf.paragraphs[0]
    p.alignment = align

    if line_gap > 0:
        pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        etree.SubElement(spcAft, qn('a:spcPts')).set('val', str(line_gap * 100))

    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(color) if not _is_rgb(color) else color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def bullet_list(slide, lines, l, t, w, h,
                font="Times New Roman", size=14, color="#2C3E50", line_gap_pt=8,
                bullet_color="#2980B9"):
    if not lines:
        return None

    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    bPr = tf._txBody.find(qn('a:bodyPr'))
    bPr.set('lIns', '91440')
    bPr.set('rIns', '0')
    bPr.set('tIns', '0')
    bPr.set('bIns', '0')

    text_rgb = _rgb(color) if isinstance(color, str) else color

    first = True
    for line in lines:
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT

        pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        etree.SubElement(spcAft, qn('a:spcPts')).set('val', str(line_gap_pt * 100))

        run = p.add_run()
        run.text = "● " + line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = text_rgb

    return tb


def render_formula_to_image(formula: str, fontsize=16) -> str:
    import tempfile, uuid, subprocess, platform, shutil, os, urllib.parse, requests

    clean_formula = formula.strip('$').strip()
    if not clean_formula:
        return None

    workdir = tempfile.mkdtemp(prefix="formula_tex_")
    try:
        latex_content = r"""\documentclass[preview,border=2pt,12pt]{standalone}
\usepackage{amsmath}
\usepackage{amssymb}
\usepackage{xcolor}
\begin{document}
$ %s $
\end{document}
""" % clean_formula

        tex_path = os.path.join(workdir, "formula.tex")
        with open(tex_path, "w", encoding="utf-8") as f:
            f.write(latex_content)

        pdflatex_cmd = shutil.which("pdflatex")
        pdftoppm_cmd = shutil.which("pdftoppm")

        if not pdflatex_cmd and platform.system() == "Windows":
            texlive_bin = r"D:\ProgramFiles\texlive\2024\bin\windows"
            if os.path.exists(os.path.join(texlive_bin, "pdflatex.exe")):
                pdflatex_cmd = os.path.join(texlive_bin, "pdflatex.exe")
                pdftoppm_cmd = os.path.join(texlive_bin, "pdftoppm.exe")

        if pdflatex_cmd and pdftoppm_cmd:
            subprocess.run(
                [pdflatex_cmd, "-interaction=batchmode", "-output-directory", workdir, tex_path],
                capture_output=True, timeout=15
            )

            pdf_path = os.path.join(workdir, "formula.pdf")
            if os.path.exists(pdf_path):
                png_base = os.path.join(workdir, "formula")
                subprocess.run(
                    [pdftoppm_cmd, "-png", "-r", "300", "-singlefile", pdf_path, png_base],
                    capture_output=True, timeout=15
                )

                png_path = png_base + ".png"
                if os.path.exists(png_path):
                    final_path = os.path.join(tempfile.gettempdir(), f'formula_local_{uuid.uuid4().hex[:8]}.png')
                    shutil.copy2(png_path, final_path)
                    print(f"    [本地渲染] 公式成功: {clean_formula[:20]}...")
                    return final_path
                else:
                    print("    [本地渲染警告] pdftoppm 转换 PNG 失败")
            else:
                print("    [本地渲染警告] pdflatex 生成 PDF 失败")
        else:
            print("    [提示] 未在系统检测到本地 LaTeX 环境。")

    except Exception as e:
        print(f"    [本地渲染异常]: {e}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

    try:
        encoded_formula = urllib.parse.quote(clean_formula)
        url = f"https://latex.codecogs.com/png.image?\\dpi{{300}}\\huge {encoded_formula}"
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(url, headers=headers, timeout=10)

        if resp.status_code == 200:
            final_path = os.path.join(tempfile.gettempdir(), f'formula_online_{uuid.uuid4().hex[:8]}.png')
            with open(final_path, 'wb') as f:
                f.write(resp.content)
            print(f"    [在线渲染] 公式成功: {clean_formula[:20]}...")
            return final_path
        else:
            print(f"    [在线渲染失败] 状态码: {resp.status_code}")
    except Exception as e:
        print(f"    [在线渲染异常]: {e}")

    return None


def is_color_dark(bg_color) -> bool:
    rgb = _rgb(bg_color)
    r, g, b = rgb[0], rgb[1], rgb[2]
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return luminance < 140


def extract_and_render_formulas(text: str, x: float, y: float, slide, theme: ThemeStyle, max_width: float = Inches(8.0), text_color=None) -> float:
    import re
    if text_color is None:
        text_color = ACADEMIC_BLUE

    pure_formula_match = re.fullmatch(r'\s*\$\$(.*?)\$\$\s*', text, flags=re.DOTALL)

    if pure_formula_match:
        formula = pure_formula_match.group(1).strip()
        formula_img = render_formula_to_image(formula)

        if formula_img and os.path.exists(formula_img):
            try:
                pic = slide.shapes.add_picture(formula_img, x, y)
                max_h = Inches(0.8)
                if pic.height > max_h:
                    ratio = max_h / pic.height
                    pic.height = int(max_h)
                    pic.width = int(pic.width * ratio)
                if pic.width > max_width:
                    ratio = max_width / pic.width
                    pic.width = int(max_width)
                    pic.height = int(pic.height * ratio)
                return y + pic.height + Inches(0.1)
            except Exception as e:
                print(f"    [渲染警告] 图片插入失败: {e}")

    clean_text = text.replace('$$', '')
    chars_per_line = 70
    estimated_lines = max(1, len(clean_text) // chars_per_line + clean_text.count('\n') + 1)
    line_height = Inches(0.35) * estimated_lines + Inches(0.05)

    textbox(slide, clean_text, x, y, max_width, line_height,
           font=theme.font_body, size=14, color=text_color,
           wrap=True, line_gap=4)

    return y + line_height


def _draw_header(slide, t: ThemeStyle, title: str, section_num: str = ""):
    rect(slide, 0, 0, W, Inches(0.04), ACCENT_BLUE)
    rect(slide, 0, Inches(0.04), W, Inches(0.95), _rgb(t.header_bg))

    if section_num:
        textbox(slide, section_num, Inches(0.3), Inches(0.08), Inches(0.8), Inches(0.85),
                font=t.font_title, size=32, color=ACCENT_BLUE, bold=True)
        title_x = Inches(1.1)
        title_w = Inches(11.7)
    else:
        title_x = Inches(0.4)
        title_w = Inches(12.5)

    textbox(slide, title, title_x, Inches(0.12), title_w, Inches(0.82),
            font=t.font_title, size=26, color=_rgb(t.header_text), bold=True)


def _draw_footer(slide, t: ThemeStyle, page_num: int = None):
    footer_bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, H - Inches(0.4), W, Inches(0.4), footer_bg)
    rect(slide, 0, H - Inches(0.4), W, Inches(0.02), ACCENT_BLUE)


def draw_cover(slide, t: ThemeStyle, title, subtitle="", author="", affiliation="", date=""):
    rect(slide, 0, 0, W, H, _rgb(t.bg_cover))
    rect(slide, 0, 0, W, Inches(0.06), ACCENT_BLUE)
    rect(slide, 0, H - Inches(2.2), W, Inches(2.2), DARK_BG)
    rect(slide, 0, 0, Inches(0.15), H, ACCENT_BLUE)

    if subtitle:
        textbox(slide, subtitle, Inches(0.6), Inches(1.8), Inches(11), Inches(0.5),
                font=t.font_body, size=14, color=ACCENT_TEAL, italic=True)

    textbox(slide, title, Inches(0.6), Inches(2.5), Inches(11.5), Inches(2.2),
            font="Times New Roman", size=40, color=WHITE, bold=True, wrap=True)

    rect(slide, Inches(0.6), Inches(5.0), Inches(3), Inches(0.03), ACCENT_BLUE)

    if author:
        textbox(slide, author, Inches(0.6), Inches(5.2), Inches(11), Inches(0.5),
                font=t.font_body, size=18, color=WHITE)
    if affiliation:
        textbox(slide, affiliation, Inches(0.6), Inches(5.8), Inches(11), Inches(0.5),
                font=t.font_body, size=14, color=RGBColor(0xBD, 0xC3, 0xC7), italic=True)
    if date:
        textbox(slide, date, Inches(0.6), Inches(6.5), Inches(11), Inches(0.4),
                font=t.font_body, size=12, color=RGBColor(0x95, 0xA5, 0xA6))


def draw_toc(slide, t: ThemeStyle, title, sections):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)
    rect(slide, 0, 0, Inches(0.12), H, ACCENT_BLUE)
    rect(slide, 0, 0, W, Inches(0.05), ACCENT_BLUE)

    textbox(slide, title, Inches(0.5), Inches(0.4), Inches(12), Inches(0.9),
            font="Times New Roman", size=32, color=ACADEMIC_BLUE, bold=True)
    rect(slide, Inches(0.5), Inches(1.25), Inches(2), Inches(0.04), ACCENT_BLUE)

    n = len(sections)
    start_y = Inches(1.8)
    step_y = Inches(0.95) if n <= 6 else Inches(0.8)

    for i, s in enumerate(sections[:8]):
        y = start_y + i * step_y
        circle_bg = ACCENT_BLUE if i % 2 == 0 else ACCENT_TEAL
        rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.4), Inches(0.4), circle_bg)
        textbox(slide, str(i + 1), Inches(0.5), y + Inches(0.1), Inches(0.4), Inches(0.4),
                font=t.font_title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, s, Inches(1.1), y, Inches(11), Inches(0.65),
                font="Times New Roman", size=18, color=ACADEMIC_BLUE)


def draw_content(slide, t: ThemeStyle, title, section_num, bullets, stats=None):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)
    _draw_header(slide, t, title, section_num)

    is_dark_bg = is_color_dark(bg)
    slide_text_color = WHITE if is_dark_bg else ACADEMIC_BLUE

    has_stats = stats and len(stats) > 0
    content_x = Inches(0.5)
    content_w = Inches(8.5) if has_stats else Inches(12.0)

    n = len(bullets)
    if n <= 5: step, start_y = Inches(0.95), Inches(1.5)
    elif n <= 8: step, start_y = Inches(0.75), Inches(1.4)
    else: step, start_y = Inches(0.65), Inches(1.3)

    current_y = start_y

    for i, b in enumerate(bullets):
        if not b: continue
        dot_bg = ACCENT_BLUE if i % 2 == 0 else ACCENT_TEAL
        rect(slide, content_x, current_y + Inches(0.12), Inches(0.12), Inches(0.12), dot_bg)

        if '$$' in b:
            current_y = extract_and_render_formulas(
                b, content_x + Inches(0.25), current_y, slide, t, content_w, text_color=slide_text_color
            )
            current_y += Inches(0.1)
        else:
            textbox(slide, b, content_x + Inches(0.25), current_y, content_w, step - Inches(0.1),
                    font="Times New Roman", size=14, color=slide_text_color, wrap=True, line_gap=4)
            current_y += step

    if has_stats:
        for i, (num, desc) in enumerate(stats[:3]):
            x = Inches(9.3)
            y = Inches(1.5) + i * Inches(1.85)

            if is_dark_bg:
                card_bg_color = RGBColor(0x2A, 0x3B, 0x4C)
                card_line = "#3A4B5C"
                desc_color = RGBColor(0xA0, 0xB0, 0xC0)
            else:
                card_bg_color = RGBColor(0xF8, 0xF9, 0xFA)
                card_line = "#DEE2E6"
                desc_color = RGBColor(0x5D, 0x6D, 0x7E)

            rect(slide, x, y, Inches(3.6), Inches(1.7), card_bg_color, line_color=card_line, line_width=0.5)
            rect(slide, x, y, Inches(0.08), Inches(1.7), ACCENT_BLUE)

            num_str = str(num)
            is_formula = "$" in num_str or "\\" in num_str

            if is_formula:
                formula_img = render_formula_to_image(num_str)
                if formula_img and os.path.exists(formula_img):
                    slide.shapes.add_picture(formula_img, x + Inches(0.15), y + Inches(0.2), width=Inches(3.3), height=Inches(0.9))
                else:
                    textbox(slide, num_str.strip('$').strip(), x + Inches(0.15), y + Inches(0.2), Inches(3.3), Inches(0.9),
                            font="Times New Roman", size=16, color=slide_text_color, align=PP_ALIGN.CENTER)
            else:
                textbox(slide, num_str, x + Inches(0.15), y + Inches(0.2), Inches(3.3), Inches(0.8),
                        font="Times New Roman", size=28, color=slide_text_color, bold=True, align=PP_ALIGN.CENTER)

            textbox(slide, desc, x + Inches(0.15), y + Inches(1.0), Inches(3.3), Inches(0.55),
                    font=t.font_body, size=11, color=desc_color, align=PP_ALIGN.CENTER, wrap=True)

    _draw_footer(slide, t)


def draw_section_divider(slide, t: ThemeStyle, num, title):
    rect(slide, 0, 0, W, H, _rgb(t.bg_cover))
    rect(slide, 0, 0, W, Inches(0.06), ACCENT_BLUE)
    rect(slide, 0, 0, Inches(0.15), H, ACCENT_BLUE)
    rect(slide, 0, H - Inches(2.0), W, Inches(2.0), DARK_BG)

    textbox(slide, num, Inches(0.6), Inches(1.5), Inches(1.5), Inches(1.5),
            font="Times New Roman", size=72, color=ACCENT_BLUE, bold=True)
    rect(slide, Inches(0.6), Inches(3.3), Inches(4), Inches(0.04), ACCENT_BLUE)
    textbox(slide, title, Inches(0.6), Inches(3.6), Inches(11), Inches(2.0),
            font="Times New Roman", size=36, color=WHITE, bold=True, wrap=True)


def draw_two_column(slide, t: ThemeStyle, title, section_num, left_title, left_bullets, right_title, right_bullets):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)
    _draw_header(slide, t, title, section_num)

    is_dark_bg = is_color_dark(bg)
    slide_text_color = WHITE if is_dark_bg else ACADEMIC_BLUE
    card_bg = RGBColor(0x2A, 0x3B, 0x4C) if is_dark_bg else (RGBColor(0xF8, 0xF9, 0xFA) if t.card_bg == "#FFFFFF" else _rgb(t.card_bg))
    card_line = "#3A4B5C" if is_dark_bg else "#E8E8E8"

    col_w, gap = Inches(6.0), Inches(0.65)
    lx = Inches(0.45)
    rx = lx + col_w + gap
    start_y = Inches(1.3)

    columns = [(lx, left_title, left_bullets), (rx, right_title, right_bullets)]
    accent_colors = [ACCENT_BLUE, ACCENT_TEAL]

    for idx, (x, col_t, col_b) in enumerate(columns):
        rect(slide, x, start_y, col_w, Inches(0.5), accent_colors[idx])
        textbox(slide, col_t, x + Inches(0.15), start_y + Inches(0.05), col_w - Inches(0.2), Inches(0.42),
                font="Times New Roman", size=15, color=WHITE, bold=True)

        n = len(col_b)
        fs = max(12, min(14, 15 - max(0, n - 5)))
        lg = max(6, int(45 / n)) if n else 8

        rect(slide, x, start_y + Inches(0.55), col_w, Inches(5.0), card_bg, line_color=card_line, line_width=0.3)
        bullet_list(slide, col_b, x + Inches(0.15), start_y + Inches(0.65), col_w - Inches(0.25), Inches(4.8),
                    font="Times New Roman", size=fs, color=slide_text_color, line_gap_pt=lg, bullet_color=accent_colors[idx])
    _draw_footer(slide, t)


def draw_image(slide, t: ThemeStyle, title, section_num, image_path, caption=""):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)
    _draw_header(slide, t, title, section_num)

    if image_path and os.path.exists(image_path):
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            max_w, max_h = Inches(11.5), Inches(4.8)
            if aspect > (max_w / max_h):
                pic_w = max_w
                pic_h = pic_w / aspect
            else:
                pic_h = max_h
                pic_w = pic_h * aspect

            pic_l = (W - pic_w) / 2
            pic_t = Inches(1.2) + (max_h - pic_h) / 2

            rect(slide, pic_l - Inches(0.05), pic_t - Inches(0.05),
                 pic_w + Inches(0.1), pic_h + Inches(0.1),
                 WHITE, line_color="#CCCCCC", line_width=0.5)

            slide.shapes.add_picture(image_path, pic_l, pic_t, width=pic_w, height=pic_h)
        except Exception as e:
            print(f"    插入图片失败: {e}")

    fig_y = H - Inches(0.9)
    rect(slide, Inches(0.5), fig_y, W - Inches(1.0), Inches(0.02), ACCENT_BLUE)
    if caption:
        cap_color = _rgb(t.text_secondary) if not t.bg_dark else GRAY_TEXT
        textbox(slide, f"图: {caption}", Inches(0.5), fig_y + Inches(0.1),
                W - Inches(1.0), Inches(0.6),
                font="Times New Roman", size=12, color=cap_color,
                italic=True, align=PP_ALIGN.CENTER)

    _draw_footer(slide, t)


def draw_table(slide, t: ThemeStyle, title, section_num, table_content, caption=""):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)
    _draw_header(slide, t, title, section_num)

    is_dark_bg = is_color_dark(bg)
    slide_text_color = WHITE if is_dark_bg else ACADEMIC_BLUE

    rows_data = []
    for line in table_content.strip().split('\n'):
        if line.strip():
            cells = [c.strip() for c in (line.split('|') if '|' in line else line.split('\t')) if c.strip()]
            if cells: rows_data.append(cells)

    if rows_data:
        num_rows = min(len(rows_data), 15)
        num_cols = max(len(row) for row in rows_data[:num_rows]) if rows_data else 3
        table_w = min(Inches(12), W - Inches(1.0))
        table_h = Inches(min(num_rows * 0.4, 4.5))
        table_l = (W - table_w) / 2
        table_t = Inches(1.3)

        tbl = slide.shapes.add_table(num_rows, num_cols, table_l, table_t, table_w, table_h).table
        for i in range(num_cols): tbl.columns[i].width = int(table_w / num_cols)

        for r, row in enumerate(rows_data[:num_rows]):
            for c, cell in enumerate(row[:num_cols]):
                cell_obj = tbl.cell(r, c)
                cell_obj.text = cell
                tf = cell_obj.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                run.text = cell
                run.font.name = "Times New Roman"
                run.font.size = Pt(11)

                if r == 0:
                    cell_obj.fill.solid()
                    cell_obj.fill.fore_color.rgb = ACCENT_BLUE
                    run.font.color.rgb = WHITE
                    run.font.bold = True
                else:
                    cell_obj.fill.solid()
                    if is_dark_bg:
                        cell_obj.fill.fore_color.rgb = RGBColor(0x2A, 0x3B, 0x4C) if r % 2 == 0 else RGBColor(0x1A, 0x25, 0x3A)
                        run.font.color.rgb = WHITE
                    else:
                        cell_obj.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA) if r % 2 == 0 else WHITE
                        run.font.color.rgb = ACADEMIC_BLUE
    else:
        textbox(slide, table_content[:500], Inches(0.5), Inches(1.3), W - Inches(1.0), Inches(4.5),
                font="Times New Roman", size=12, color=slide_text_color, wrap=True)

    if caption:
        cap_y = H - Inches(0.8)
        rect(slide, Inches(0.5), cap_y, W - Inches(1.0), Inches(0.02), ACCENT_BLUE)
        cap_color = RGBColor(0xA0, 0xB0, 0xC0) if is_dark_bg else _rgb(t.text_secondary)
        textbox(slide, f"表: {caption}", Inches(0.5), cap_y + Inches(0.08), W - Inches(1.0), Inches(0.5),
                font="Times New Roman", size=12, color=cap_color, italic=True, align=PP_ALIGN.CENTER)
    _draw_footer(slide, t)


def draw_conclusion(slide, t: ThemeStyle, title, contributions, future_work, closing="感谢聆听！欢迎提问与交流"):
    bg = _rgb(t.bg_content) if not t.bg_dark else _rgb(t.bg_cover)
    rect(slide, 0, 0, W, H, bg)

    is_dark_bg = is_color_dark(bg)
    slide_text_color = WHITE if is_dark_bg else ACADEMIC_BLUE
    card_bg = RGBColor(0x2A, 0x3B, 0x4C) if is_dark_bg else (RGBColor(0xF8, 0xF9, 0xFA) if t.card_bg == "#FFFFFF" else _rgb(t.card_bg))
    card_line = "#3A4B5C" if is_dark_bg else "#E8E8E8"

    rect(slide, 0, 0, W, Inches(0.05), ACCENT_BLUE)
    rect(slide, 0, Inches(0.05), W, Inches(0.95), _rgb(t.header_bg))
    textbox(slide, title, Inches(0.5), Inches(0.12), Inches(12), Inches(0.82), font="Times New Roman", size=28, color=_rgb(t.header_text), bold=True)

    contrib_x, contrib_w = Inches(0.5), Inches(6.2)
    rect(slide, contrib_x, Inches(1.2), contrib_w, Inches(5.3), card_bg, line_color=card_line, line_width=0.3)
    rect(slide, contrib_x, Inches(1.2), contrib_w, Inches(0.5), ACCENT_BLUE)
    textbox(slide, "主要贡献", contrib_x + Inches(0.2), Inches(1.25), contrib_w - Inches(0.4), Inches(0.45), font="Times New Roman", size=16, color=WHITE, bold=True)

    n_c = len(contributions)
    step_c = min(Inches(0.95), Inches(4.5) / max(n_c, 1))

    for i, c in enumerate(contributions[:5]):
        if not c: continue
        cy = Inches(1.85) + i * step_c
        circle_bg = ACCENT_BLUE if i % 2 == 0 else ACCENT_TEAL
        rect(slide, contrib_x + Inches(0.2), cy, Inches(0.35), Inches(0.35), circle_bg)
        textbox(slide, str(i + 1), contrib_x + Inches(0.2), cy, Inches(0.35), Inches(0.35), font="Times New Roman", size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, c, contrib_x + Inches(0.7), cy - Inches(0.05), contrib_w - Inches(0.9), step_c - Inches(0.1), font="Times New Roman", size=13, color=slide_text_color, wrap=True)

    future_x, future_w = Inches(7.0), Inches(5.8)
    rect(slide, future_x, Inches(1.2), future_w, Inches(5.3), card_bg, line_color=card_line, line_width=0.3)
    rect(slide, future_x, Inches(1.2), future_w, Inches(0.5), ACCENT_TEAL)
    textbox(slide, "未来展望", future_x + Inches(0.2), Inches(1.25), future_w - Inches(0.4), Inches(0.45), font="Times New Roman", size=16, color=WHITE, bold=True)

    n_f = len(future_work)
    step_f = min(Inches(0.9), Inches(4.5) / max(n_f, 1))

    for i, f in enumerate(future_work[:5]):
        if not f: continue
        fy = Inches(1.85) + i * step_f
        rect(slide, future_x + Inches(0.2), fy + Inches(0.08), Inches(0.25), Inches(0.2), ACCENT_TEAL)
        textbox(slide, f, future_x + Inches(0.6), fy - Inches(0.05), future_w - Inches(0.8), step_f - Inches(0.1), font="Times New Roman", size=13, color=slide_text_color, wrap=True)

    rect(slide, 0, H - Inches(0.95), W, Inches(0.95), ACADEMIC_BLUE)
    textbox(slide, closing, Inches(0.5), H - Inches(0.9), W - Inches(1.0), Inches(0.9), font="Times New Roman", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

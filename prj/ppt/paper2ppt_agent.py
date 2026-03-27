"""
Paper to PPT Agent v4
Pipeline: OCR → Content Extraction → Vision Analysis → LLM Outline → PPT Rendering
"""
import os, re, json, base64, requests, tempfile, shutil
from typing import Dict, List, Optional
from pptx import Presentation

from ocr_extractor import call_ocr_api, extract_content as ocr_extract_content
from ppt_style_analyzer import ThemeStyle, analyze_template_style, _fallback_theme
from ppt_renderer import (
    draw_cover, draw_toc, draw_content,
    draw_two_column, draw_section_divider, draw_conclusion,
    W, H
)

try:
    from config import get_ocr_config, get_llm_config
    _default_config = True
except ImportError:
    _default_config = False


class Paper2PPTAgent:

    def __init__(self, api_token: str,
                 ocr_api_url: str = None,
                 llm_api_url: str = None,
                 model_name: str = None,
                 vision_model_name: str = None):
        self.api_token = api_token

        if ocr_api_url:
            self.ocr_api_url = ocr_api_url
        elif _default_config:
            cfg = get_ocr_config()
            self.ocr_api_url = cfg.get("api_url") or "https://pdx5w3wdkbyap7h4.aistudio-app.com/layout-parsing"
        else:
            self.ocr_api_url = "https://pdx5w3wdkbyap7h4.aistudio-app.com/layout-parsing"

        if llm_api_url:
            self.llm_api_url = llm_api_url
        elif _default_config:
            cfg = get_llm_config()
            self.llm_api_url = cfg.get("api_url") or "https://aistudio.baidu.com/llm/lmapi/v3"
        else:
            self.llm_api_url = "https://aistudio.baidu.com/llm/lmapi/v3"

        if model_name:
            self.model_name = model_name
        elif _default_config:
            cfg = get_llm_config()
            self.model_name = cfg.get("model_name", "ernie-4.5-turbo-128k-preview")
        else:
            self.model_name = "ernie-4.5-turbo-128k-preview"

        if vision_model_name:
            self.vision_model_name = vision_model_name
        elif _default_config:
            cfg = get_llm_config()
            self.vision_model_name = cfg.get("vision_model_name", "ernie-4.5-turbo-vl")
        else:
            self.vision_model_name = "ernie-4.5-turbo-vl"

        self.templates = {
            "1": "ppts/模板1_经典深海蓝.pptx",
            "2": "ppts/模板2_现代简洁白.pptx",
            "3": "ppts/模板3_暗色科技风.pptx",
            "4": "ppts/模板4_学术典雅.pptx",
        }

    def _call_llm(self, prompt: str, temperature: float = 0.7) -> str:
        headers = {
            "Authorization": f"token {self.api_token}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.model_name,
            "messages": [
                {"role": "system", "content": "你是专业学术PPT制作专家。"},
                {"role": "user",   "content": prompt},
            ],
            "temperature": temperature,
            "max_completion_tokens": 10000,
        }
        resp = requests.post(f"{self.llm_api_url}/chat/completions", json=payload, headers=headers)
        if resp.status_code != 200:
            raise Exception(f"LLM失败: {resp.status_code}, {resp.text}")
        return resp.json()["choices"][0]["message"]["content"]

    def _call_llm_vision(self, prompt: str, image_paths: List[str]) -> str:
        content = []
        for p in image_paths[:4]:
            with open(p, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            content.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{b64}"}
            })
        content.append({"type": "text", "text": prompt})

        headers = {
            "Authorization": f"token {self.api_token}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.vision_model_name,
            "messages": [{"role": "user", "content": content}],
            "temperature": 0.2,
            "max_completion_tokens": 1000,
        }
        resp = requests.post(f"{self.llm_api_url}/chat/completions", json=payload, headers=headers)
        if resp.status_code != 200:
            raise Exception(f"视觉LLM失败: {resp.status_code}, {resp.text[:300]}")
        return resp.json()["choices"][0]["message"]["content"]

    def ocr_paper(self, file_path: str) -> Dict:
        return call_ocr_api(
            file_path,
            api_url=self.ocr_api_url,
            api_token=self.api_token,
            use_chart_recognition=False,
        )

    def extract_paper_content(self, ocr_data: Dict, output_dir: str = None) -> Dict:
        return ocr_extract_content(ocr_data, output_dir=output_dir)

    def analyze_images_with_vision(self, paper_content: Dict) -> List[str]:
        images = paper_content.get("images", [])
        if not images:
            print("  无图片可分析")
            return []

        descriptions = []
        for i, img in enumerate(images[:10]):
            local_path = img.get("local_path")
            if not local_path or not os.path.exists(local_path):
                descriptions.append(f"图片{i}: (无法访问)")
                continue

            try:
                prompt = (
                    "请用中文简要描述这张论文中的图片内容。"
                    "如果是实验结果图，说明它展示的指标和趋势。"
                    "如果是架构图，说明模型的主要组成部分。"
                    "如果是数据表格截图，说明表格内容。"
                    "请用2-3句话概括，不超过100字。"
                )
                desc = self._call_llm_vision(prompt, [local_path])
                descriptions.append(f"图片{i}({img.get('filename','')})：{desc.strip()}")
                print(f"  图片{i}分析完成: {desc.strip()[:50]}...")
            except Exception as e:
                descriptions.append(f"图片{i}: (分析失败: {e})")
                print(f"  图片{i}分析失败: {e}")

        return descriptions

    def analyze_template_visually(self, template_id: str) -> ThemeStyle:
        tpl_path = self.templates.get(template_id)
        if not tpl_path or not os.path.exists(tpl_path):
            print(f"  模板文件不存在，使用默认配色")
            return _fallback_theme(template_id)

        return analyze_template_style(
            pptx_path=tpl_path,
            llm_func=self._call_llm_vision,
            theme_id=template_id,
        )

    def generate_ppt_structure(self, paper_content: Dict, requirements: str,
                               time_limit: str, display_mode: str,
                               theme: ThemeStyle,
                               image_descriptions: List[str] = None) -> Dict:
        print("正在生成PPT大纲...")

        full_text_preview = paper_content['full_text'][:8000]
        sections_info = "\n".join([f"{i+1}. {s['title']}" for i, s in enumerate(paper_content['sections'])])
        images_info = f"可用图片数量：{len(paper_content['images'])}张"
        formulas_info = f"提取公式数量：{len(paper_content['formulas'])}个"
        tables_info = f"识别表格数量：{len(paper_content['tables'])}个"

        img_desc_block = ""
        if image_descriptions:
            img_desc_block = "\n【图片内容分析（由视觉AI生成）】\n" + "\n".join(image_descriptions)

        formula_block = ""
        if paper_content.get('formulas'):
            top_formulas = paper_content['formulas'][:15]
            formula_block = "\n【提取的重要公式】\n" + "\n".join(
                [f"  公式{i+1}: ${f}$" for i, f in enumerate(top_formulas)]
            )

        table_block = ""
        if paper_content.get('tables'):
            for i, tbl in enumerate(paper_content['tables'][:5]):
                table_block += f"\n  表格{i+1}(第{tbl['page']}页):\n{tbl['content'][:300]}\n"
            table_block = "\n【提取的表格数据】" + table_block

        max_img_idx = max(0, len(paper_content['images']) - 1)

        prompt = f"""你是专业的学术演讲PPT制作专家。现在你需要根据以下论文内容，从汇报演讲的角度生成一份详尽的PPT大纲，确保听众能够清楚理解这篇论文的核心内容。

【论文信息】
标题：{paper_content['title']}
摘要：{paper_content['abstract'][:1500]}

章节结构：
{sections_info}

{images_info}
{formulas_info}
{tables_info}
{img_desc_block}
{formula_block}
{table_block}

论文内容预览（前8000字）：
{full_text_preview}

【演讲需求】
需求：{requirements or '无特殊要求'}
时间：{time_limit}（约每页1-1.5分钟）
方式：{display_mode}

【PPT模板风格】{theme.style_notes}

【核心原则——从汇报角度思考】
1. 听众视角：PPT是给听众看的，不是照搬论文。每页内容要让没读过论文的人也能理解核心信息。
2. 讲故事：按"问题-动机-方法-结果-总结"顺序铺排，逻辑清晰循序渐进。
3. 图片优先：论文中的实验结果图、模型架构图、流程图等只要不是重复含义的都要放入PPT。
4. 公式不能少：论文中最重要的模型公式和参数公式必须展示，并在前后用 content 页解释公式含义。
5. 表格数据要展示：实验数据表、对比表、算法流程表都要放入 PPT。
6. bullets要具体：每个要点20-50字，包含具体的技术细节和数据，不要写空洞的标题。
7. 数据说话：尽量用具体数字、百分比、对比结果来说明观点。

【页数要求】
- PPT总计 20-30 页
- 每个主要章节 3-5 页内容页
- 方法论部分要最详细（5-8页）
- 实验结果要充实（4-6页，包含多张实验图/表）

【可用幻灯片类型】
- cover     封面（首页，只用1次）
- toc       目录（只用1次，自动从内容中提取章节标题）
- section   章节过渡页（每大章节前用一次，3-5次）
- content   内容页（最常用：标题+详细要点列表，每页4-6个要点）
- two_col   两栏对比页（适合方法对比、前后对比、优缺点对比）
- image     图片展示页（展示论文中的图表、架构图、实验结果图）
- table     表格展示页（展示实验数据表格、算法伪代码等）
- conclusion 结论页（末页，只用1次）

【字段说明】
cover:      title, subtitle, author, affiliation, date
toc:        title, sections（列表，4-6个主要章节标题，由你从PPT内容中提取概括）
section:    title, section_num（"01"/"02"等）
content:    title, section_num, bullets（列表，每条20-50字，含具体信息）, stats（可选，[["数字","说明"],...]最多3个）
two_col:    title, section_num, left_title, left_bullets（4-5条）, right_title, right_bullets（4-5条）
image:      title, section_num, image_index（0-{max_img_idx}）, caption（图片说明，描述图中关键信息）
table:      title, section_num, table_content（表格内容，markdown格式）, caption（表格说明）
conclusion: title, contributions（4-5条详细贡献）, future_work（3-4条未来方向）

【特别注意】
- image 类型的 image_index 范围是 0 到 {max_img_idx}，请合理分配
- 每张不同含义的图片都应该展示
- 公式应在 content 页的 bullets 中用$$包裹展示，并附带文字解释其含义
- toc 页的 sections 列表应当是你从PPT整体内容中总结出的目录标题
- 目录标题应简洁有信息量

请直接返回JSON格式的PPT结构，确保内容详细完整，页数在20-30页之间。只返回JSON，不要额外文字。"""

        response = self._call_llm(prompt, temperature=0.5)
        try:
            start_idx = -1
            for i, char in enumerate(response):
                if char in '{[':
                    start_idx = i
                    break

            end_idx = -1
            for i in range(len(response) - 1, -1, -1):
                if response[i] in '}]':
                    end_idx = i
                    break

            if start_idx != -1 and end_idx != -1:
                json_str = response[start_idx:end_idx+1]
            else:
                json_str = response

            raw_data = json.loads(json_str)

            structure = {"slides": []}
            if isinstance(raw_data, dict):
                if "slides" in raw_data:
                    slides_data = raw_data["slides"]
                    if isinstance(slides_data, list) and len(slides_data) == 1 and isinstance(slides_data[0], dict):
                        first = slides_data[0]
                        nested_keys = ["cover", "toc", "section", "content", "two_col", "image", "table", "conclusion"]
                        has_nested = any(k in first for k in nested_keys)
                        if has_nested:
                            for key in nested_keys:
                                if key in first:
                                    val = first[key]
                                    if isinstance(val, list):
                                        for item in val:
                                            if isinstance(item, dict):
                                                item["type"] = key
                                                structure["slides"].append(item)
                                    elif isinstance(val, dict):
                                        val["type"] = key
                                        structure["slides"].append(val)
                        else:
                            structure = raw_data
                    else:
                        structure = raw_data
                else:
                    structure["slides"] = [raw_data]
            elif isinstance(raw_data, list):
                for item in raw_data:
                    if isinstance(item, dict):
                        if "type" in item:
                            structure["slides"].append(item)
                        else:
                            for k, v in item.items():
                                v["type"] = k
                                structure["slides"].append(v)
                                break

            print(f"大纲生成完成，共 {len(structure.get('slides',[]))} 页")
            return structure
        except json.JSONDecodeError as e:
            print(f"JSON解析失败: {e}")
            try:
                debug_path = "output/failed_llm_response.txt"
                os.makedirs("output", exist_ok=True)
                with open(debug_path, "w", encoding="utf-8") as f:
                    f.write(response)
                print(f"原始响应已保存至: {debug_path}")
            except: pass
            raise e

    def download_image(self, url, path) -> bool:
        try:
            resp = requests.get(url, timeout=10)
            if resp.status_code == 200:
                with open(path, 'wb') as f: f.write(resp.content)
                return True
        except: pass
        return False

    def create_ppt(self, ppt_structure: Dict, paper_content: Dict,
                   theme: ThemeStyle, output_path: str) -> str:
        print(f"\n正在从零绘制PPT（风格：{theme.name}）...")

        prs = Presentation()
        prs.slide_width  = W
        prs.slide_height = H
        blank = prs.slide_layouts[6]

        temp_dir = tempfile.mkdtemp(prefix="ppt_images_")
        image_files = []
        for i, img in enumerate(paper_content.get("images", [])[:20]):
            safe_name = os.path.basename(img.get("filename", f"img_{i}.jpg"))
            p = os.path.join(temp_dir, safe_name)
            if img.get("local_path") and os.path.exists(img["local_path"]):
                shutil.copy2(img["local_path"], p)
                image_files.append(p)
                print(f"  使用本地图片 {i}: {img.get('filename', 'unknown')}")
            elif self.download_image(img.get("url"), p):
                image_files.append(p)
                print(f"  下载图片 {i}: {img.get('filename', 'unknown')}")

        section_counter = 0
        for slide_info in ppt_structure.get("slides", []):
            stype = slide_info.get("type", "content")
            slide = prs.slides.add_slide(blank)
            print(f"  [{stype:10s}] {slide_info.get('title','')[:35]}")

            try:
                if stype == "cover":
                    draw_cover(slide, theme,
                               title=slide_info.get("title", paper_content["title"]),
                               subtitle=slide_info.get("subtitle","学术报告"),
                               author=slide_info.get("author",""),
                               affiliation=slide_info.get("affiliation",""),
                               date=slide_info.get("date",""))

                elif stype == "toc":
                    draw_toc(slide, theme,
                             title=slide_info.get("title","目录"),
                             sections=slide_info.get("sections",[]))

                elif stype == "section":
                    section_counter += 1
                    draw_section_divider(slide, theme,
                                         num=slide_info.get("section_num", f"{section_counter:02d}"),
                                         title=slide_info.get("title",""))

                elif stype == "two_col":
                    draw_two_column(slide, theme,
                                    title=slide_info.get("title",""),
                                    section_num=slide_info.get("section_num",""),
                                    left_title=slide_info.get("left_title",""),
                                    left_bullets=slide_info.get("left_bullets",[]),
                                    right_title=slide_info.get("right_title",""),
                                    right_bullets=slide_info.get("right_bullets",[]))

                elif stype == "image":
                    img_idx = slide_info.get("image_index", 0)
                    img_path = image_files[img_idx] if 0 <= img_idx < len(image_files) else None
                    from ppt_renderer import draw_image
                    draw_image(slide, theme,
                              title=slide_info.get("title",""),
                              section_num=slide_info.get("section_num",""),
                              image_path=img_path,
                              caption=slide_info.get("caption",""))

                elif stype == "table":
                    from ppt_renderer import draw_table
                    draw_table(slide, theme,
                              title=slide_info.get("title",""),
                              section_num=slide_info.get("section_num",""),
                              table_content=slide_info.get("table_content",""),
                              caption=slide_info.get("caption",""))

                elif stype == "conclusion":
                    draw_conclusion(slide, theme,
                                    title=slide_info.get("title","结论与展望"),
                                    contributions=slide_info.get("contributions",[]),
                                    future_work=slide_info.get("future_work",[]))

                else:
                    stats_raw = slide_info.get("stats")
                    stats = [(str(s[0]), str(s[1])) for s in stats_raw if len(s) >= 2] if stats_raw else None
                    draw_content(slide, theme,
                                 title=slide_info.get("title",""),
                                 section_num=slide_info.get("section_num",""),
                                 bullets=slide_info.get("bullets",[]),
                                 stats=stats)

            except Exception as e:
                print(f"    [Error] {e}")
                import traceback; traceback.print_exc()

        prs.save(output_path)
        print(f"\n✅ PPT 已保存: {output_path}  (共 {len(prs.slides)} 页)")
        shutil.rmtree(temp_dir, ignore_errors=True)
        return output_path

    def process(self, paper_path: str, requirements: str = "",
                time_limit: str = "15分钟", display_mode: str = "线下演讲",
                template_id: str = "1", output_path: str = "output.pptx") -> str:
        sep = "=" * 60
        out_dir = os.path.dirname(output_path) or "output"
        os.makedirs(out_dir, exist_ok=True)
        try:
            print(f"\n{sep}\n步骤 1/4: OCR识别\n{sep}")
            ocr_result = self.ocr_paper(paper_path)

            print(f"\n{sep}\n步骤 2/4: 提取论文内容\n{sep}")
            paper_content = self.extract_paper_content(ocr_result, output_dir=out_dir)

            print(f"\n{sep}\n步骤 3/5: 视觉分析图片内容\n{sep}")
            image_descriptions = self.analyze_images_with_vision(paper_content)

            print(f"\n{sep}\n步骤 4/5: 应用模板配色\n{sep}")
            theme = _fallback_theme(template_id)
            print(f"  配色方案: bg={theme.bg_cover}, accent={theme.accent}, "
                  f"header={theme.header_bg} ({theme.name})")

            print(f"\n{sep}\n步骤 5/5: 生成大纲 + 从零绘制PPT\n{sep}")
            ppt_structure = self.generate_ppt_structure(
                paper_content, requirements, time_limit, display_mode, theme,
                image_descriptions=image_descriptions
            )

            debug_dir = os.path.dirname(output_path) or "."
            os.makedirs(debug_dir, exist_ok=True)
            with open(os.path.join(debug_dir, "debug_paper_content.json"), "w", encoding="utf-8") as f:
                json.dump(paper_content, f, ensure_ascii=False, indent=2)
            with open(os.path.join(debug_dir, "debug_ppt_structure.json"), "w", encoding="utf-8") as f:
                json.dump(ppt_structure, f, ensure_ascii=False, indent=2)

            result = self.create_ppt(ppt_structure, paper_content, theme, output_path)

            print(f"\n{sep}\n✅ 完成！→ {result}\n{sep}")
            return result

        except Exception as e:
            print(f"\n{sep}\n❌ 失败: {e}\n{sep}")
            import traceback; traceback.print_exc()
            raise


if __name__ == "__main__":
    print("请通过 app.py 或 workflows/paper2ppt.py 启动本模块")
